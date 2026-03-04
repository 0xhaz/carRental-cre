#!/usr/bin/env python3
"""
Generate RegShield Presentation Slides (.pptx)
Run: python3 generate-presentation.py
Output: RegShield_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

# ── Color Palette ───────────────────────────────────────────────────────────

DARK_BG = RGBColor(15, 23, 42)        # slate-900
BLUE_PRIMARY = RGBColor(59, 130, 246)  # blue-500
BLUE_DARK = RGBColor(30, 64, 175)     # blue-800
GREEN = RGBColor(34, 197, 94)         # green-500
PURPLE = RGBColor(147, 51, 234)       # purple-600
ORANGE = RGBColor(249, 115, 22)       # orange-500
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(203, 213, 225)  # slate-300
MID_GRAY = RGBColor(148, 163, 184)    # slate-400
DARK_GRAY = RGBColor(51, 65, 85)      # slate-700
CHAINLINK_BLUE = RGBColor(55, 91, 210)
WORLDID_BLACK = RGBColor(0, 0, 0)
CARD_BG = RGBColor(30, 41, 59)        # slate-800

# ── Helpers ─────────────────────────────────────────────────────────────────

def add_bg(slide, color=DARK_BG):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, radius=None):
    """Add a rounded rectangle background shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=BLUE_PRIMARY):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.space_before = Pt(4)
        # Bullet character
        run = p.add_run()
        run.text = "\u25B8 "  # small triangle bullet
        run.font.size = Pt(font_size)
        run.font.color.rgb = bullet_color
        run.font.name = "Calibri"
        # Item text
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox

def add_card(slide, left, top, width, height, title, body_items,
             accent_color=BLUE_PRIMARY, title_size=16, body_size=13):
    """Add a card with title and bullet items."""
    # Card background
    card = add_shape_bg(slide, left, top, width, height, CARD_BG, radius=0.05)
    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(0.15),
                 width - Inches(0.4), Inches(0.4),
                 title, font_size=title_size, color=WHITE, bold=True)
    # Body items
    y = top + Inches(0.6)
    for item in body_items:
        add_text_box(slide, left + Inches(0.35), y,
                     width - Inches(0.5), Inches(0.35),
                     f"\u2022 {item}", font_size=body_size, color=LIGHT_GRAY)
        y += Inches(0.32)

def add_number_stat(slide, left, top, number, label, color=BLUE_PRIMARY):
    """Add a big number with label below."""
    add_text_box(slide, left, top, Inches(2.2), Inches(0.7),
                 number, font_size=36, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.6), Inches(2.2), Inches(0.4),
                 label, font_size=13, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)

def slide_number(slide, num):
    """Add slide number at bottom right."""
    add_text_box(slide, Inches(12.5), Inches(7.0), Inches(0.7), Inches(0.3),
                 str(num), font_size=10, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Gradient accent line at top
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(13.333), Inches(0.06))
line.fill.solid()
line.fill.fore_color.rgb = BLUE_PRIMARY
line.line.fill.background()

# Title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "RegShield", font_size=60, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
             "Tokenized Vehicle Rental Platform", font_size=28,
             color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.6),
             "Invest in Vehicles. Earn Real Revenue. All On-Chain.",
             font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Sponsor badges
for i, (label, color) in enumerate([
    ("Chainlink CRE", CHAINLINK_BLUE),
    ("WorldID", PURPLE),
    ("ERC-3643", GREEN),
]):
    x = Inches(3.8) + Inches(i * 2.1)
    badge = add_shape_bg(slide, x, Inches(4.8), Inches(1.8), Inches(0.45), color, radius=0.15)
    # Add text to shape
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(2)

# Date
add_text_box(slide, Inches(4), Inches(6.2), Inches(5), Inches(0.4),
             datetime.date.today().strftime("Chainlink Block Magic Hackathon — %B %Y"),
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

slide_number(slide, 1)

# Speaker notes
notes = slide.notes_slide.notes_text_frame
notes.text = (
    "Hi everyone, I'm [Your Name] and today I'm presenting RegShield — "
    "a tokenized vehicle rental platform built for the Chainlink Block Magic Hackathon.\n\n"
    "RegShield lets investors buy fractional ownership of real vehicles using blockchain, "
    "and earn real rental revenue — all enforced on-chain with ERC-3643 compliance tokens.\n\n"
    "We heavily leverage two sponsors: Chainlink CRE for decentralized automation, "
    "and WorldID for sybil-resistant investor onboarding. Let me walk you through how it all works."
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "The Problem", font_size=36, color=WHITE, bold=True)

# Problem cards
problems = [
    ("Centralized Backend", [
        "Cron jobs for booking, campaigns, revenue sync",
        "Single point of failure — no verifiability",
        "Admin manually approves KYC & milestones",
    ], ORANGE),
    ("No On-Chain Compliance", [
        "Vehicle investments lack regulatory framework",
        "No automated identity verification",
        "Token transfers without compliance checks",
    ], RGBColor(239, 68, 68)),
    ("Sybil Vulnerability", [
        "One person can create multiple investor accounts",
        "No proof-of-personhood verification",
        "Gaming of investment allocations",
    ], PURPLE),
    ("Manual Milestone Verification", [
        "Admin manually checks VIN, insurance, registration",
        "Escrow funds released without automated verification",
        "No real-time vehicle telemetry monitoring",
    ], MID_GRAY),
]

for i, (title, items, color) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.1)
    y = Inches(1.7) + Inches(row * 2.8)
    add_card(slide, x, y, Inches(5.7), Inches(2.5), title, items, color)

slide_number(slide, 2)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    "So what problems are we solving?\n\n"
    "First — centralized backend. Traditional rental platforms rely on cron jobs running on a single server "
    "for bookings, campaigns, and revenue. If that server goes down, everything stops. There's no verifiability.\n\n"
    "Second — no on-chain compliance. Vehicle investments today have no regulatory framework enforced by code. "
    "Token transfers happen without any compliance checks.\n\n"
    "Third — sybil vulnerability. Without proof-of-personhood, one person can create multiple investor accounts "
    "and game investment allocations.\n\n"
    "Fourth — manual milestone verification. An admin has to manually check VIN numbers, insurance documents, "
    "and registration before releasing escrow funds. This is slow, error-prone, and doesn't scale.\n\n"
    "RegShield solves ALL of these with Chainlink CRE and WorldID."
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3: OUR SOLUTION
# ═══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Our Solution: RegShield", font_size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "A fully on-chain compliant vehicle investment platform with decentralized automation",
             font_size=16, color=MID_GRAY)

# Solution pillars
pillars = [
    ("Chainlink CRE", CHAINLINK_BLUE, [
        "5 off-chain workflows replace backend cron jobs",
        "Automated VIN validation via NHTSA API",
        "Milestone-based escrow verification",
        "Vehicle telematics & compliance monitoring",
    ]),
    ("WorldID", PURPLE, [
        "Proof-of-personhood for investor onboarding",
        "Orb-level ZK verification — sybil resistant",
        "On-chain via WorldIDVerifier contract",
        "One human = one verified investor account",
    ]),
    ("ERC-3643 Tokens", GREEN, [
        "AssetToken — fractional vehicle ownership",
        "RevenueToken — quarterly rental income rights",
        "Identity Registry enforces compliance",
        "Transfer restrictions by investor type & KYC",
    ]),
]

for i, (title, color, items) in enumerate(pillars):
    x = Inches(0.6) + Inches(i * 4.2)
    add_card(slide, x, Inches(2.2), Inches(3.9), Inches(4.2),
             title, items, color, title_size=20, body_size=14)

slide_number(slide, 3)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    "Our solution has three pillars.\n\n"
    "Chainlink CRE — this is the big one. We built 5 off-chain workflows that completely replace "
    "our backend cron jobs. These run on Chainlink's Decentralized Oracle Network, so there's no single point of failure. "
    "They handle everything from automated VIN validation via the NHTSA API, to milestone-based escrow verification, "
    "to vehicle telematics monitoring.\n\n"
    "WorldID — we use Orb-level zero-knowledge verification for investor onboarding. "
    "This means one human equals one verified investor account. No sybil attacks possible. "
    "The verification happens through our WorldIDVerifier smart contract on Sepolia.\n\n"
    "ERC-3643 Security Tokens — every vehicle gets two compliance tokens. "
    "An AssetToken representing fractional ownership, and a RevenueToken for quarterly rental income rights. "
    "The Identity Registry enforces transfer restrictions based on investor type and KYC status."
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4: HOW IT WORKS
# ═══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "How It Works", font_size=36, color=WHITE, bold=True)

steps = [
    ("1", "Invest", "Investor selects a vehicle\ncampaign and commits ETH.\nFunds locked in PaymentEscrow.", BLUE_PRIMARY),
    ("2", "CRE Verifies", "Chainlink CRE auto-verifies\n4 milestones: VIN, purchase,\ninsurance, registration.", CHAINLINK_BLUE),
    ("3", "Tokens Minted", "ERC-3643 AssetTokens &\nRevenueTokens minted\nto verified investors.", GREEN),
    ("4", "Earn Revenue", "Vehicle generates rental\nincome. Quarterly waterfall:\n50% investors, 10% operator.", ORANGE),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5) + Inches(i * 3.2)
    y = Inches(1.8)

    # Step number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(1.1), y, Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)

    # Arrow between steps (except last)
    if i < 3:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(2.7), y + Inches(0.2),
            Inches(0.7), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK_GRAY
        arrow.line.fill.background()

    # Title
    add_text_box(slide, x, y + Inches(1.1), Inches(3), Inches(0.5),
                 title, font_size=20, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, x, y + Inches(1.6), Inches(3), Inches(1.5),
                 desc, font_size=13, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

# Bottom highlight box
highlight = add_shape_bg(slide, Inches(1), Inches(5.5),
                         Inches(11.3), Inches(1.2), CARD_BG, radius=0.03)
add_text_box(slide, Inches(1.3), Inches(5.6), Inches(10.7), Inches(0.4),
             "Chainlink CRE replaces 4 centralized backend services + adds 1 new capability",
             font_size=16, color=CHAINLINK_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.3), Inches(6.0), Inches(10.7), Inches(0.5),
             "Onboarding  |  Campaign Monitor  |  Rental Lifecycle  |  Vehicle Telematics  |  Compliance Monitor",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

slide_number(slide, 4)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    "Here's the user journey in four steps.\n\n"
    "Step 1 — An investor selects a vehicle campaign on our marketplace and commits ETH. "
    "The funds are locked in our PaymentEscrow smart contract. Nothing gets released until milestones are verified.\n\n"
    "Step 2 — This is where Chainlink CRE comes in. Our CRE workflows automatically verify 4 milestones: "
    "VIN validation through the NHTSA API, proof of purchase, insurance verification, and vehicle registration. "
    "No manual admin approval needed.\n\n"
    "Step 3 — Once milestones are verified, ERC-3643 AssetTokens and RevenueTokens are minted "
    "to the verified investors. These are compliance tokens — only KYC-verified, WorldID-verified investors can hold them.\n\n"
    "Step 4 — The vehicle goes live on our rental platform and generates income. "
    "Revenue is distributed quarterly through a waterfall model: 50% to investors, 10% to the vehicle operator, "
    "and the rest covers maintenance and platform fees.\n\n"
    "The key insight here — Chainlink CRE replaces 4 centralized backend services and adds 1 entirely new capability "
    "that wasn't possible before: real-time compliance monitoring."
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5: CHAINLINK CRE DEEP DIVE
# ═══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Chainlink CRE Integration", font_size=36, color=CHAINLINK_BLUE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "5 Decentralized Workflows Replace Centralized Backend Cron Jobs",
             font_size=16, color=MID_GRAY)

# Workflow table
workflows = [
    ("Onboarding", "Auto-verify investor identity (ERC-3643 + WorldID), approve/reject", "Manual admin approval", "5 reports"),
    ("Campaign Monitor", "Detect expired/underfunded campaigns, trigger batch refunds", "campaignScheduler.js", "2 reports"),
    ("Rental Lifecycle", "4-milestone verification: VIN, escrow, insurance, registration", "bookingScheduler.js", "3 reports"),
    ("Vehicle Telematics", "Poll odometer/telemetry data, update VehicleNFT on-chain", "revenueSyncService.js", "3 reports"),
    ("Compliance", "Monitor registration/insurance expiry, overdue maintenance", "None (NEW capability)", "3 reports"),
]

# Table header
header_y = Inches(1.9)
header_bg = add_shape_bg(slide, Inches(0.6), header_y, Inches(12.1), Inches(0.5), BLUE_DARK, radius=0.02)
cols = [("CRE Workflow", 2.2), ("What It Does", 4.8), ("Replaces", 2.6), ("Output", 1.5)]
x = Inches(0.8)
for label, w in cols:
    add_text_box(slide, x, header_y + Inches(0.05), Inches(w), Inches(0.4),
                 label, font_size=13, color=WHITE, bold=True)
    x += Inches(w + 0.2)

# Table rows
for i, (wf, desc, replaces, output) in enumerate(workflows):
    y = Inches(2.5) + Inches(i * 0.85)
    row_color = CARD_BG if i % 2 == 0 else DARK_BG
    row_bg = add_shape_bg(slide, Inches(0.6), y, Inches(12.1), Inches(0.75), row_color, radius=0.01)

    x = Inches(0.8)
    values = [(wf, 2.2, True, BLUE_PRIMARY), (desc, 4.8, False, LIGHT_GRAY),
              (replaces, 2.6, False, ORANGE if "NEW" in replaces else MID_GRAY),
              (output, 1.5, False, GREEN)]
    for val, w, b, c in values:
        add_text_box(slide, x, y + Inches(0.1), Inches(w), Inches(0.55),
                     val, font_size=12, color=c, bold=b)
        x += Inches(w + 0.2)

# Footer
add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
             "All 5 workflows compile to WASM, run on Chainlink DON nodes, and submit reports to 5 Receiver contracts on Sepolia",
             font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

slide_number(slide, 5)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    "Let's dive deeper into our Chainlink CRE integration — this is the core of our project.\n\n"
    "We built 5 TypeScript workflows that compile to WASM and run on Chainlink DON nodes.\n\n"
    "The Onboarding workflow replaces manual admin approval. It auto-verifies investor identity "
    "by checking ERC-3643 Identity Registry and WorldID status, then generates 5 onboarding reports.\n\n"
    "The Campaign Monitor workflow replaces our campaignScheduler.js cron job. It detects expired or underfunded "
    "campaigns and triggers batch refunds automatically.\n\n"
    "The Rental Lifecycle workflow replaces bookingScheduler.js. It handles 4-milestone verification: "
    "VIN lookup, escrow check, insurance validation, and registration confirmation.\n\n"
    "The Vehicle Telematics workflow replaces revenueSyncService.js. It polls odometer and telemetry data "
    "and updates the VehicleNFT on-chain.\n\n"
    "And finally, the Compliance Monitor — this is a completely NEW capability that didn't exist before. "
    "It monitors registration expiry, insurance expiry, and overdue maintenance in real-time. "
    "This was impossible with our centralized backend.\n\n"
    "All 5 workflows submit ABI-encoded reports to 5 Receiver contracts deployed on Sepolia."
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6: WORLDID INTEGRATION
# ═══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "WorldID Integration", font_size=36, color=PURPLE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "Proof-of-Personhood for Sybil-Resistant Investor Onboarding",
             font_size=16, color=MID_GRAY)

# Flow steps
flow_steps = [
    ("1. Verify in\nWorld App", "Investor opens RegShield\nin World App. Completes\nOrb-level ZK verification."),
    ("2. Cloud Proof\nValidation", "Backend API verifies proof\nvia Worldcoin cloud service.\nReturns merkle root + nullifier."),
    ("3. On-Chain\nRegistration", "WorldIDVerifier.sol calls\nWorld ID Router. Maps address\nas verified. Prevents duplicates."),
    ("4. Compliance\nGate", "useComplianceStatus() hook\nchecks isWorldIDVerified()\nbefore enabling investments."),
]

for i, (title, desc) in enumerate(flow_steps):
    x = Inches(0.5) + Inches(i * 3.2)
    y = Inches(2.0)

    add_card(slide, x, y, Inches(2.9), Inches(3.0),
             title, [], PURPLE, title_size=15)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.9),
                 Inches(2.4), Inches(1.8),
                 desc, font_size=12, color=LIGHT_GRAY)

    # Connector arrow
    if i < 3:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(2.95), y + Inches(1.2),
            Inches(0.3), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = PURPLE
        arrow.line.fill.background()

# Smart contract highlight
sc_box = add_shape_bg(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.5), CARD_BG, radius=0.03)

add_text_box(slide, Inches(1.1), Inches(5.5), Inches(5), Inches(0.4),
             "WorldIDVerifier.sol — Sepolia", font_size=16, color=PURPLE, bold=True)

contract_details = [
    "Address: 0x838C9397F5c00f7010924dDc4c2E93Fcab6c0363",
    "verifyAndRegister(signal, root, nullifierHash, proof) — registers verified user",
    "isWorldIDVerified(address) — quick lookup for compliance gates",
    "Nullifier hash prevents same human from verifying twice (sybil resistance)",
]
add_bullet_list(slide, Inches(1.1), Inches(5.9), Inches(10.5), Inches(1.0),
                contract_details, font_size=11, color=LIGHT_GRAY, bullet_color=PURPLE)

slide_number(slide, 6)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    "Now let's look at our WorldID integration — our second sponsor technology.\n\n"
    "The flow has 4 steps.\n\n"
    "First, the investor opens RegShield in the World App and completes Orb-level ZK verification. "
    "This generates a zero-knowledge proof that proves they're a unique human without revealing their identity.\n\n"
    "Second, our backend API validates the proof via the Worldcoin cloud verification service. "
    "It returns a merkle root and nullifier hash.\n\n"
    "Third, our WorldIDVerifier smart contract — deployed at the address shown below — "
    "calls the World ID Router contract to verify the proof on-chain. "
    "It maps the wallet address as verified and stores the nullifier hash to prevent the same human from verifying twice.\n\n"
    "Fourth, our frontend has a useComplianceStatus() hook that checks isWorldIDVerified() "
    "before enabling any investment actions. If you're not WorldID verified, you simply cannot invest.\n\n"
    "The key function is verifyAndRegister — it takes the signal, root, nullifier hash, and proof, "
    "validates everything on-chain, and registers the user. The nullifier hash is the sybil resistance mechanism — "
    "same human, same nullifier, can't register twice."
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7: ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Architecture Overview", font_size=36, color=WHITE, bold=True)

# Three layers
layers = [
    ("Frontend (Next.js)", BLUE_PRIMARY, Inches(0.4), [
        "Wagmi + Viem for wallet & contract interaction",
        "Chainlink Price Feed hook (ETH/USD live)",
        "WorldID MiniKit integration",
        "ERC-3643 compliance status hooks",
        "Role-based dashboards (Investor, Rentor, Renter, Admin)",
    ]),
    ("Smart Contracts (Sepolia)", GREEN, Inches(2.8), [
        "34 contracts: VehicleNFT, PaymentEscrow, IdentityRegistry",
        "ERC-3643 AssetToken & RevenueToken factories",
        "WorldIDVerifier — on-chain proof-of-personhood",
        "5 CRE Receiver contracts (MockKeystoneForwarder)",
        "RevenueDistributor — waterfall: 50% investors, 10% operator",
    ]),
    ("Chainlink CRE (Off-Chain)", CHAINLINK_BLUE, Inches(5.2), [
        "5 TypeScript workflows compiled to WASM",
        "Run on DON nodes with consensus aggregation",
        "Read on-chain state + call external APIs (NHTSA)",
        "Submit ABI-encoded reports to Receiver contracts",
        "Replace: campaignScheduler, bookingScheduler, revenueSyncService",
    ]),
]

for title, color, y_offset, items in layers:
    y = Inches(1.4) + y_offset
    # Layer band
    band = add_shape_bg(slide, Inches(0.6), y, Inches(12.1), Inches(2.0), CARD_BG, radius=0.03)
    # Color indicator
    indicator = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.1), Inches(2.0))
    indicator.fill.solid()
    indicator.fill.fore_color.rgb = color
    indicator.line.fill.background()
    # Title
    add_text_box(slide, Inches(1.0), y + Inches(0.15), Inches(4), Inches(0.4),
                 title, font_size=18, color=color, bold=True)
    # Items split into two columns
    mid = len(items) // 2 + len(items) % 2
    add_bullet_list(slide, Inches(1.0), y + Inches(0.55), Inches(5.5), Inches(1.3),
                    items[:mid], font_size=12, color=LIGHT_GRAY, bullet_color=color)
    add_bullet_list(slide, Inches(6.8), y + Inches(0.55), Inches(5.5), Inches(1.3),
                    items[mid:], font_size=12, color=LIGHT_GRAY, bullet_color=color)

slide_number(slide, 7)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    "Here's our three-layer architecture.\n\n"
    "The Frontend is built with Next.js. We use Wagmi and Viem for wallet and contract interaction, "
    "Chainlink Price Feed hooks for live ETH/USD prices, WorldID MiniKit for identity verification, "
    "and custom ERC-3643 compliance status hooks. We have role-based dashboards for all four user types.\n\n"
    "The Smart Contract layer has 34 contracts deployed on Sepolia. Key ones include VehicleNFT for the vehicle registry, "
    "PaymentEscrow for milestone-locked funds, the ERC-3643 token factories for AssetTokens and RevenueTokens, "
    "WorldIDVerifier for proof-of-personhood, and 5 CRE Receiver contracts that accept reports from the DON.\n\n"
    "The Chainlink CRE layer is the off-chain compute. Our 5 TypeScript workflows compile to WASM "
    "and run on DON nodes with consensus aggregation. They can read on-chain state and call external APIs "
    "like NHTSA for VIN validation. They submit ABI-encoded reports back to the Receiver contracts.\n\n"
    "The important thing to note — CRE replaces our backend cron jobs: campaignScheduler, bookingScheduler, "
    "and revenueSyncService. These were single points of failure. Now they run decentralized."
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8: SMART CONTRACTS
# ═══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Smart Contracts on Sepolia", font_size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
             "34 contracts deployed — key contracts shown below",
             font_size=14, color=MID_GRAY)

contract_groups = [
    ("Identity & Compliance", GREEN, [
        ("IdentityRegistry", "0x188C...927A", "ERC-3643 on-chain identity"),
        ("WorldIDVerifier", "0x838C...0363", "Proof-of-personhood (WorldID)"),
        ("ComplianceRules", "0xe81a...2C0f", "Transfer restriction enforcement"),
    ]),
    ("Tokens & Revenue", BLUE_PRIMARY, [
        ("AssetTokenFactory", "0x1EEE...38ef", "ERC-3643 ownership tokens"),
        ("RevenueTokenFactory", "0x1eDE...206C", "ERC-3643 income rights"),
        ("RevenueDistributor", "0xD08b...538D", "Waterfall distribution engine"),
    ]),
    ("Vehicle & Payments", ORANGE, [
        ("VehicleNFT", "0xfcE0...1814", "ERC-721 vehicle registry"),
        ("InvestmentEscrow", "0x8e4d...bc89", "Milestone-locked escrow"),
        ("RentalPaymentProtocol", "0xE7a7...fc2c", "Rental payment processing"),
    ]),
    ("CRE Receivers", CHAINLINK_BLUE, [
        ("OnboardingReceiver", "0xF080...3850", "Identity verification reports"),
        ("VehicleReceiver", "0x73c5...e81B", "Telematics & mileage reports"),
        ("PaymentReceiver", "0x2f7F...B7C2", "Milestone completion reports"),
    ]),
]

for i, (group_name, color, contracts) in enumerate(contract_groups):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + Inches(col * 6.3)
    y = Inches(1.8) + Inches(row * 2.8)

    add_text_box(slide, x, y, Inches(5.8), Inches(0.4),
                 group_name, font_size=16, color=color, bold=True)

    for j, (name, addr, desc) in enumerate(contracts):
        cy = y + Inches(0.45) + Inches(j * 0.65)
        row_bg = add_shape_bg(slide, x, cy, Inches(5.8), Inches(0.55), CARD_BG, radius=0.02)
        add_text_box(slide, x + Inches(0.15), cy + Inches(0.02),
                     Inches(2.2), Inches(0.25),
                     name, font_size=12, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(2.4), cy + Inches(0.02),
                     Inches(1.5), Inches(0.25),
                     addr, font_size=10, color=MID_GRAY, font_name="Courier New")
        add_text_box(slide, x + Inches(0.15), cy + Inches(0.28),
                     Inches(5.5), Inches(0.25),
                     desc, font_size=10, color=LIGHT_GRAY)

slide_number(slide, 8)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    "We have 34 smart contracts deployed on Sepolia. Let me highlight the key ones.\n\n"
    "In the Identity and Compliance group — the IdentityRegistry is the ERC-3643 on-chain identity store. "
    "The WorldIDVerifier handles proof-of-personhood verification. ComplianceRules enforces transfer restrictions.\n\n"
    "In Tokens and Revenue — AssetTokenFactory and RevenueTokenFactory deploy ERC-3643 compliance tokens "
    "for each vehicle. The RevenueDistributor handles the waterfall distribution engine — "
    "it automatically splits rental revenue: 50% to investors, 10% to the operator.\n\n"
    "In Vehicle and Payments — VehicleNFT is our ERC-721 vehicle registry where each car is an NFT. "
    "InvestmentEscrow holds funds locked until milestones are verified. "
    "RentalPaymentProtocol processes rental payments.\n\n"
    "And our CRE Receivers — we have 5 Receiver contracts, one for each CRE workflow. "
    "OnboardingReceiver accepts identity verification reports, VehicleReceiver accepts telematics data, "
    "and PaymentReceiver accepts milestone completion reports. "
    "All of these are verified on Sepolia and use the MockKeystoneForwarder for simulation."
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9: PLATFORM ROLES
# ═══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Platform Roles", font_size=36, color=WHITE, bold=True)

roles = [
    ("Investor", BLUE_PRIMARY, [
        "Browse marketplace & invest ETH",
        "WorldID verification for sybil resistance",
        "KYC upload → ERC-3643 identity on-chain",
        "Receive AssetTokens + RevenueTokens",
        "Claim quarterly rental revenue",
    ]),
    ("Rentor", GREEN, [
        "Register vehicles as ERC-721 NFTs",
        "Create fundraising campaigns",
        "Manage rental bookings",
        "Earn 10% operator fee on revenue",
        "Track analytics & vehicle performance",
    ]),
    ("Renter", PURPLE, [
        "Browse & filter available vehicles",
        "Book vehicles with date selection",
        "Optional WorldID verification",
        "Leave reviews & ratings",
        "Track active & past rentals",
    ]),
    ("Admin", ORANGE, [
        "Review & approve KYC submissions",
        "Monitor CRE activity feed",
        "Track milestone verification status",
        "Manage disputes & refunds",
        "Revenue distribution oversight",
    ]),
]

for i, (role, color, features) in enumerate(roles):
    x = Inches(0.4) + Inches(i * 3.2)
    add_card(slide, x, Inches(1.6), Inches(3.0), Inches(5.2),
             role, features, color, title_size=22, body_size=13)

slide_number(slide, 9)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    "RegShield supports four user roles.\n\n"
    "Investors can browse the vehicle marketplace, invest ETH into campaigns, and earn quarterly rental revenue. "
    "They must complete WorldID verification for sybil resistance and upload KYC documents "
    "which get registered on-chain via ERC-3643. Once verified, they receive AssetTokens and RevenueTokens.\n\n"
    "Rentors — that's vehicle owners — can register their vehicles as ERC-721 NFTs, "
    "create fundraising campaigns to tokenize their vehicles, manage rental bookings, "
    "and earn a 10% operator fee on all revenue generated.\n\n"
    "Renters can browse and filter available vehicles, book them with date selection, "
    "optionally verify with WorldID for trusted rentals, and leave reviews.\n\n"
    "Admins review KYC submissions, monitor the CRE activity feed to see all automated actions, "
    "track milestone verification status, and manage disputes and refunds.\n\n"
    "Each role has its own dedicated dashboard in the frontend with real-time data."
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10: CRE DEMO RESULTS
# ═══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "CRE Simulation Results", font_size=36, color=CHAINLINK_BLUE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "All 5 workflows pass simulation — ready for DON deployment",
             font_size=16, color=MID_GRAY)

# Results
results = [
    ("onboarding", "5 onboarding reports", "\u2705"),
    ("campaign", "2 campaign reports", "\u2705"),
    ("rental", "3 milestone reports", "\u2705"),
    ("vehicle", "3 telemetry reports", "\u2705"),
    ("compliance", "3 compliance reports", "\u2705"),
]

for i, (wf, result, status) in enumerate(results):
    y = Inches(2.0) + Inches(i * 0.9)
    row_bg = add_shape_bg(slide, Inches(1), y, Inches(11.3), Inches(0.75), CARD_BG, radius=0.02)

    add_text_box(slide, Inches(1.3), y + Inches(0.15), Inches(0.5), Inches(0.45),
                 status, font_size=20)
    add_text_box(slide, Inches(2.0), y + Inches(0.15), Inches(3), Inches(0.45),
                 f"{wf}-workflow", font_size=16, color=WHITE, bold=True,
                 font_name="Courier New")
    add_text_box(slide, Inches(5.5), y + Inches(0.15), Inches(4), Inches(0.45),
                 f"[SIMULATION] {result}", font_size=14, color=GREEN)

# Stats bar
stat_y = Inches(6.6)
add_number_stat(slide, Inches(1.5), stat_y - Inches(0.3), "5", "CRE Workflows", CHAINLINK_BLUE)
add_number_stat(slide, Inches(4.0), stat_y - Inches(0.3), "16", "Reports Generated", GREEN)
add_number_stat(slide, Inches(6.5), stat_y - Inches(0.3), "5", "Receiver Contracts", ORANGE)
add_number_stat(slide, Inches(9.0), stat_y - Inches(0.3), "34", "Smart Contracts", PURPLE)

slide_number(slide, 10)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    "Here are our CRE simulation results — all 5 workflows pass successfully.\n\n"
    "The onboarding workflow generates 5 reports covering investor identity verification. "
    "The campaign workflow produces 2 reports for expired campaign detection. "
    "The rental workflow generates 3 milestone verification reports. "
    "The vehicle workflow produces 3 telemetry reports with odometer data. "
    "And the compliance workflow generates 3 compliance monitoring reports.\n\n"
    "That's 16 total reports across 5 workflows, all submitted to 5 Receiver contracts on Sepolia.\n\n"
    "These workflows are currently running in simulation mode using the CRE SDK simulator. "
    "The simulator validates the WASM binaries, consensus logic, and report encoding. "
    "All 5 are ready for DON deployment once we have access to a live Chainlink node.\n\n"
    "We have 34 smart contracts deployed on Sepolia backing all of this."
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11: TECH STACK
# ═══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Tech Stack", font_size=36, color=WHITE, bold=True)

stacks = [
    ("Frontend", BLUE_PRIMARY, [
        "Next.js 14 (App Router)",
        "TypeScript + TailwindCSS",
        "Wagmi v3 + Viem v2",
        "Thirdweb SDK v5",
        "Framer Motion",
    ]),
    ("Smart Contracts", GREEN, [
        "Solidity 0.8.24",
        "Foundry (Forge + Cast)",
        "34 contracts on Sepolia",
        "ERC-721, ERC-3643",
        "OpenZeppelin v5",
    ]),
    ("Backend", ORANGE, [
        "Node.js + Express",
        "MongoDB + Mongoose",
        "JWT Authentication",
        "node-cron schedulers",
        "ethers.js v6",
    ]),
    ("Chainlink CRE", CHAINLINK_BLUE, [
        "TypeScript workflows",
        "CRE SDK + Bun compiler",
        "WASM compilation target",
        "DON consensus aggregation",
        "5 Receiver contracts",
    ]),
    ("Identity", PURPLE, [
        "WorldID MiniKit v1.9",
        "Orb-level ZK proofs",
        "OnchainID (T-REX)",
        "ERC-3643 IdentityRegistry",
        "KYC document system",
    ]),
]

for i, (name, color, items) in enumerate(stacks):
    x = Inches(0.3) + Inches(i * 2.6)
    add_card(slide, x, Inches(1.6), Inches(2.4), Inches(4.8),
             name, items, color, title_size=16, body_size=12)

slide_number(slide, 11)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    "Quick overview of our tech stack.\n\n"
    "Frontend — Next.js 14 with App Router, TypeScript, TailwindCSS for styling, "
    "Wagmi v3 and Viem v2 for blockchain interaction, Thirdweb SDK for wallet connectivity, "
    "and Framer Motion for animations.\n\n"
    "Smart Contracts — Solidity 0.8.24, built and tested with Foundry. "
    "34 contracts on Sepolia including ERC-721, ERC-3643 tokens, and OpenZeppelin v5.\n\n"
    "Backend — Node.js with Express, MongoDB with Mongoose for data storage, "
    "JWT authentication, node-cron schedulers for the centralized services that CRE replaces, "
    "and ethers.js v6 for blockchain interaction.\n\n"
    "Chainlink CRE — TypeScript workflows compiled to WASM using the CRE SDK and Bun, "
    "running on DON nodes with consensus aggregation, with 5 Receiver contracts.\n\n"
    "Identity — WorldID MiniKit v1.9 for Orb-level ZK proofs, "
    "OnchainID from the T-REX protocol for ERC-3643 identity, "
    "and a full KYC document management system."
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12: CLOSING / THANK YOU
# ═══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(13.333), Inches(0.06))
line.fill.solid()
line.fill.fore_color.rgb = BLUE_PRIMARY
line.line.fill.background()

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "Thank You", font_size=54, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(2.8), Inches(9), Inches(0.7),
             "RegShield — Tokenized Vehicle Rental Platform",
             font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key takeaways
takeaways = [
    "Chainlink CRE replaces 4 backend services + adds 1 new capability",
    "WorldID provides sybil-resistant proof-of-personhood for investors",
    "ERC-3643 security tokens enforce on-chain regulatory compliance",
    "34 smart contracts deployed on Sepolia testnet",
    "5 CRE workflows — all passing simulation, ready for DON deployment",
]
add_bullet_list(slide, Inches(3), Inches(3.8), Inches(7), Inches(2.5),
                takeaways, font_size=16, color=LIGHT_GRAY, bullet_color=BLUE_PRIMARY)

# Sponsor badges (larger)
for i, (label, color) in enumerate([
    ("Powered by Chainlink CRE", CHAINLINK_BLUE),
    ("Secured by WorldID", PURPLE),
    ("Compliant via ERC-3643", GREEN),
]):
    x = Inches(2.5) + Inches(i * 3.0)
    badge = add_shape_bg(slide, x, Inches(6.3), Inches(2.7), Inches(0.5), color, radius=0.15)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(3)

slide_number(slide, 12)

notes = slide.notes_slide.notes_text_frame
notes.text = (
    "To wrap up — here are our key takeaways.\n\n"
    "Chainlink CRE replaces 4 centralized backend services and adds 1 entirely new capability — "
    "real-time compliance monitoring — that wasn't possible with our traditional architecture.\n\n"
    "WorldID provides sybil-resistant proof-of-personhood, ensuring one human equals one verified investor. "
    "No fake accounts, no gaming the system.\n\n"
    "ERC-3643 security tokens enforce regulatory compliance on-chain. "
    "Every token transfer is checked against the Identity Registry and compliance rules.\n\n"
    "We deployed 34 smart contracts on Sepolia and built 5 CRE workflows — "
    "all passing simulation and ready for DON deployment.\n\n"
    "RegShield demonstrates how Chainlink CRE and WorldID can transform "
    "a traditional rental platform into a fully decentralized, compliant, and automated system.\n\n"
    "Thank you for your time. I'm happy to answer any questions."
)


# ── Save ────────────────────────────────────────────────────────────────────

output_path = "/Volumes/extreme/Projects/Solidity/Portfolio/RegShield/rental-cre/demo/RegShield_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
